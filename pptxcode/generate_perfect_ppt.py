import os
import sys

# 確保在 Windows 環境下輸出 UTF-8 編碼以防止 Emoji 導致 cp950 編碼出錯
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8')

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ==============================================================================
# 🎨 核心設計系統與色彩常數 (極致高質感商務美學)
# ==============================================================================
COLOR_DARK_BG = RGBColor(15, 23, 42)       # 深藍灰色 (#0F172A)
COLOR_LIGHT_BG = RGBColor(248, 250, 252)   # 淺藍灰色 (#F8FAFC)
COLOR_PRIMARY_DARK = RGBColor(30, 41, 59)  # 深石板灰 (#1E293B)
COLOR_GOLD = RGBColor(212, 163, 89)        # 香檳金 (#D4A359)
COLOR_WHITE = RGBColor(255, 255, 255)      # 純白
COLOR_TEXT_MUTED = RGBColor(100, 116, 139)  # 灰色文字 (#64748B)
COLOR_BORDER_GRAY = RGBColor(226, 232, 240) # 淺灰邊框 (#E2E8F0)
COLOR_ROW_ALT = RGBColor(241, 245, 249)     # 交錯行淺灰色 (#F1F5F9)
COLOR_GREEN = RGBColor(16, 185, 129)       # 綠色指標 (#10B981)
COLOR_RED = RGBColor(239, 68, 68)          # 紅色指標 (#EF4444)

# ==============================================================================
# 🛠️ 輔助排版工具函數
# ==============================================================================
def set_slide_background(slide, color):
    """設置投影片單色背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, subtitle_text, dark_mode=False):
    """添加投影片統一頁首"""
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.13), Inches(1.1))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    
    # 主標題
    p_title = tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "微軟正黑體"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_GOLD if dark_mode else COLOR_PRIMARY_DARK
    p_title.space_after = Pt(2)
    
    # 副標題
    p_sub = tf.add_paragraph()
    p_sub.text = subtitle_text
    p_sub.font.name = "微軟正黑體"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = COLOR_GOLD if dark_mode else COLOR_TEXT_MUTED

def add_card(slide, x, y, w, h, fill_color, border_color=None):
    """添加圓角矩形卡片作為資訊背板"""
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_styled_table(slide, x, y, w, h, headers, rows, col_widths=None):
    """建立並格式化精美表格"""
    rows_count = len(rows) + 1
    cols_count = len(headers)
    table_shape = slide.shapes.add_table(rows_count, cols_count, x, y, w, h)
    table = table_shape.table
    
    # 設置欄寬
    if col_widths:
        for idx, width in enumerate(col_widths):
            if idx < cols_count:
                table.columns[idx].width = width

    # 格式化表頭
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY_DARK
        cell.vertical_anchor = 3 # Middle
        
        # 邊框設定 (簡化為無邊框，利用底色分隔)
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "微軟正黑體"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD
        
    # 格式化資料列
    for row_idx, row_data in enumerate(rows):
        # 決定此列底色 (交錯底色)
        row_bg = COLOR_ROW_ALT if row_idx % 2 == 1 else COLOR_WHITE
        
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            cell.vertical_anchor = 3 # Middle
            
            p = cell.text_frame.paragraphs[0]
            p.text = str(text)
            p.font.name = "微軟正黑體"
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_PRIMARY_DARK
            
            # 對齊方式：首欄靠左，數據欄靠右
            if col_idx == 0:
                p.alignment = PP_ALIGN.LEFT
                p.font.bold = True
            else:
                p.alignment = PP_ALIGN.RIGHT
                
            # 針對增減變動資料特別上色 (綠色為成長，紅色為衰退)
            if col_idx == 3 and isinstance(text, str):
                if text.startswith("+"):
                    p.font.color.rgb = COLOR_GREEN
                    p.font.bold = True
                elif text.startswith("-"):
                    p.font.color.rgb = COLOR_RED
                    p.font.bold = True
                    
    return table_shape

# ==============================================================================
# 🎯 主程式：生成精美簡報
# ==============================================================================
def generate_presentation():
    print("🚀 開始生成 臺灣證券交易所 113 年簡明財務報告簡報...")
    
    prs = Presentation()
    # 設置投影片比例為 16:9 寬螢幕
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # --------------------------------------------------------------------------
    # 📌 Slide 1: 封面 (深色高質感商務美學)
    # --------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_DARK_BG)
    
    # 黃金裝飾框
    decor = slide1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.73), Inches(5.9))
    decor.fill.background()
    decor.line.color.rgb = COLOR_GOLD
    decor.line.width = Pt(1.5)
    
    # 標題文字方塊
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.93), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    # 頂部小標
    p0 = tf.paragraphs[0]
    p0.text = "TAIWAN STOCK EXCHANGE CORPORATION"
    p0.alignment = PP_ALIGN.CENTER
    p0.font.name = "微軟正黑體"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD
    p0.space_after = Pt(20)
    
    # 主標題
    p1 = tf.add_paragraph()
    p1.text = "113 年度簡明財務報告"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "微軟正黑體"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(10)
    
    # 英文副標題
    p2 = tf.add_paragraph()
    p2.text = "Concise Financial Report 2024"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "微軟正黑體"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_GOLD
    p2.space_after = Pt(40)
    
    # 日期與發布單位
    p3 = tf.add_paragraph()
    p3.text = "中華民國 115 年 5 月發布  |  合併暨個體財務報告分析"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = "微軟正黑體"
    p3.font.size = Pt(12)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    
    # --------------------------------------------------------------------------
    # 📌 Slide 2: 大綱與核心亮點 (淺色圓角卡片)
    # --------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_LIGHT_BG)
    add_header(slide2, "簡報大綱 & 核心財務亮點", "Agenda & Key Financial Highlights")
    
    # 設計三個欄位卡片
    card_w = Inches(3.77)
    card_h = Inches(4.3)
    card_y = Inches(1.8)
    card_gap = Inches(0.4)
    
    highlights = [
        {
            "num": "01",
            "title": "合併財務狀況",
            "subtitle": "Consolidated Financial Position",
            "desc": "臺灣證券交易所合併財務體系規模擴大，合併資產總額與權益穩健成長。\n\n• 不動產及設備略降 2.37%\n• 其他資產成長達 16.85%\n• 合併權益總額穩定攀升 12.82%"
        },
        {
            "num": "02",
            "title": "個體財務狀況",
            "subtitle": "Parent Company Financial Position",
            "desc": "本體業務與投資成果豐碩，流動性與財務結構安全無虞。\n\n• 個體資產規模大幅擴張\n• 個體保留盈餘顯著增長\n• 資本公積與股東權益健康"
        },
        {
            "num": "03",
            "title": "獲利表現亮眼",
            "subtitle": "Outstanding Profitability",
            "desc": "113 年度獲利動能強勁，無論合併或個體皆創下亮眼成績。\n\n• 合併營業利益增長 55.81%\n• 合併稅前淨利增長 48.94%\n• 每股盈餘 (EPS) 達 10.17 元"
        }
    ]
    
    for idx, item in enumerate(highlights):
        x = Inches(0.6) + idx * (card_w + card_gap)
        # 背景卡片
        add_card(slide2, x, card_y, card_w, card_h, COLOR_WHITE, COLOR_BORDER_GRAY)
        
        # 內容文字方塊
        tb = slide2.shapes.add_textbox(x + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        # 序號
        p_num = tf.paragraphs[0]
        p_num.text = item["num"]
        p_num.font.name = "微軟正黑體"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_GOLD
        p_num.space_after = Pt(10)
        
        # 標題
        p_title = tf.add_paragraph()
        p_title.text = item["title"]
        p_title.font.name = "微軟正黑體"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_DARK
        
        # 英文標題
        p_sub = tf.add_paragraph()
        p_sub.text = item["subtitle"]
        p_sub.font.name = "微軟正黑體"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED
        p_sub.space_after = Pt(14)
        
        # 描述
        p_desc = tf.add_paragraph()
        p_desc.text = item["desc"]
        p_desc.font.name = "微軟正黑體"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_PRIMARY_DARK
        p_desc.line_spacing = 1.3
        
    # --------------------------------------------------------------------------
    # 📌 Slide 3: 合併簡明資產負債表 (精美原生表格)
    # --------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_LIGHT_BG)
    add_header(slide3, "合併簡明資產負債表分析", "Consolidated Concise Balance Sheet (單位：新台幣仟元)")
    
    # 準備表格數據
    headers = ["會計項目", "113 年年底 (2024)", "112 年年底 (2023)", "增減金額 / 百分比"]
    rows = [
        ["不動產及設備 Property & Equipment", "2,872,155", "2,941,919", "-69,764 (-2.37%)"],
        ["其他資產 Other Assets", "48,158,344", "41,214,567", "+6,943,777 (+16.85%)"],
        ["流動負債 (分配前) Current Liabilities", "65,360,048", "53,358,322", "+12,001,726 (+22.49%)"],
        ["非控制權益 Non-controlling Interests", "24,023,423", "21,004,655", "+3,018,768 (+14.37%)"],
        ["權益總額 (分配前) Total Equity", "112,495,787", "99,713,660", "+12,782,127 (+12.82%)"],
        ["權益總額 (分配後) Total Equity", "–", "96,121,777", "–"],
        ["保留盈餘 (分配前) Retained Earnings", "70,861,609", "62,598,326", "+8,263,283 (+13.20%)"],
        ["保留盈餘 (分配後) Retained Earnings", "–", "59,006,443", "–"]
    ]
    
    # 設定各欄寬度
    widths = [Inches(4.5), Inches(2.4), Inches(2.4), Inches(2.83)]
    
    # 背景白背板
    add_card(slide3, Inches(0.6), Inches(1.6), Inches(12.13), Inches(5.3), COLOR_WHITE, COLOR_BORDER_GRAY)
    
    # 表格位置
    add_styled_table(slide3, Inches(0.6), Inches(1.6), Inches(12.13), Inches(4.5), headers, rows, widths)
    
    # 底部註記
    tb_note = slide3.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.4))
    p_note = tb_note.text_frame.paragraphs[0]
    p_note.text = "💡 分析重點：113 年底合併權益總額(分配前)達 1,124.9 億元，成長 12.82%；流動負債與非控制權益亦皆有兩位數成長，顯示交易所集團規模與財務韌性同步加強。"
    p_note.font.name = "微軟正黑體"
    p_note.font.size = Pt(11)
    p_note.font.color.rgb = COLOR_TEXT_MUTED
    p_note.font.italic = True

    # --------------------------------------------------------------------------
    # 📌 Slide 4: 合併簡明綜合損益表 (左表右圖)
    # --------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_LIGHT_BG)
    add_header(slide4, "合併簡明綜合損益表與成長指標", "Consolidated Comprehensive Income & Growth Chart")
    
    # 左側：數據白背板與表格
    left_x = Inches(0.6)
    left_w = Inches(6.2)
    add_card(slide4, left_x, Inches(1.7), left_w, Inches(5.1), COLOR_WHITE, COLOR_BORDER_GRAY)
    
    # 左側表格數據
    headers_inc = ["獲利指標", "113 年度", "112 年度", "增減百分比"]
    rows_inc = [
        ["營業利益 Operating Profit", "19,323,103", "12,401,644", "+55.81%"],
        ["稅前淨利 Profit Before Tax", "21,733,674", "14,592,473", "+48.94%"],
        ["歸屬母公司淨利 Owners Net Profit", "13,554,170", "9,100,923", "+48.93%"],
        ["每股盈餘 (元) EPS (NT$)", "10.17 元", "6.83 元", "+48.90%"],
        ["其他綜合利益淨額 OCI, net", "(197,042)", "374,771", "-152.58%"],
        ["歸屬母公司綜合利益", "13,354,886", "9,409,857", "+41.92%"]
    ]
    widths_inc = [Inches(2.5), Inches(1.2), Inches(1.2), Inches(1.3)]
    add_styled_table(slide4, left_x, Inches(1.7), left_w, Inches(3.6), headers_inc, rows_inc, widths_inc)
    
    # 右側：原生數據柱狀圖
    right_x = Inches(7.1)
    right_w = Inches(5.6)
    add_card(slide4, right_x, Inches(1.7), right_w, Inches(5.1), COLOR_WHITE, COLOR_BORDER_GRAY)
    
    # 建立圖表數據 (億元為單位，讓圖表數值清爽)
    chart_data = CategoryChartData()
    chart_data.categories = ["營業利益", "稅前淨利", "歸屬母公司淨利"]
    chart_data.add_series("112 年度", (124.0, 145.9, 91.0))
    chart_data.add_series("113 年度", (193.2, 217.3, 135.5))
    
    # 繪製圖表
    chart_shape = slide4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        right_x + Inches(0.4), Inches(2.2), right_w - Inches(0.8), Inches(4.2), chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.font.name = "微軟正黑體"
    chart.legend.font.size = Pt(10)
    chart.has_title = True
    chart.chart_title.text_frame.text = "關鍵獲利指標對比 (新台幣 億元)"
    chart.chart_title.text_frame.paragraphs[0].font.name = "微軟正黑體"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY_DARK
    
    # --------------------------------------------------------------------------
    # 📌 Slide 5: 個體簡明資產負債表 (精美原生表格)
    # --------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_LIGHT_BG)
    add_header(slide5, "個體簡明資產負債表分析 (臺灣證券交易所本體)", "Parent Company Only Concise Balance Sheet (單位：新台幣仟元)")
    
    headers_ind = ["會計項目", "113 年年底 (2024)", "112 年年底 (2023)", "增減金額 / 百分比"]
    rows_ind = [
        ["不動產及設備 Property & Equipment", "1,966,301", "2,040,821", "-74,520 (-3.65%)"],
        ["其他資產 Other Assets", "50,938,082", "44,552,720", "+6,385,362 (+14.33%)"],
        ["流動負債 (分配前) Current Liabilities", "61,571,110", "49,542,382", "+12,028,728 (+24.28%)"],
        ["流動負債 (分配後) Current Liabilities", "–", "53,134,265", "–"],
        ["負債總額 (分配前) Total Liabilities", "62,881,912", "50,888,749", "+11,993,163 (+23.57%)"],
        ["負債總額 (分配後) Total Liabilities", "–", "54,480,632", "–"],
        ["權益總額 (分配前) Total Equity", "88,472,364", "78,709,005", "+9,763,359 (+12.40%)"],
        ["權益總額 (分配後) Total Equity", "–", "75,117,122", "–"],
        ["保留盈餘 (分配前) Retained Earnings", "70,861,609", "62,598,326", "+8,263,283 (+13.20%)"],
        ["保留盈餘 (分配後) Retained Earnings", "–", "59,006,443", "–"],
        ["資本公積 Capital Surplus", "3,924", "3,568", "+356 (+9.98%)"],
        ["其他權益 Other Equity Interest", "4,282,105", "4,520,392", "-238,287 (-5.27%)"]
    ]
    
    widths_ind = [Inches(4.5), Inches(2.4), Inches(2.4), Inches(2.83)]
    
    # 背景白背板 (加大版面)
    add_card(slide5, Inches(0.6), Inches(1.5), Inches(12.13), Inches(5.6), COLOR_WHITE, COLOR_BORDER_GRAY)
    
    # 表格位置
    add_styled_table(slide5, Inches(0.6), Inches(1.5), Inches(12.13), Inches(5.2), headers_ind, rows_ind, widths_ind)

    # --------------------------------------------------------------------------
    # 📌 Slide 6: 個體簡明綜合損益表 (左表右圖)
    # --------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_LIGHT_BG)
    add_header(slide6, "個體簡明綜合損益表與獲利指標 (證交所本體)", "Parent Company Comprehensive Income & Profit Chart")
    
    # 左側：數據白背板與表格
    add_card(slide6, left_x, Inches(1.7), left_w, Inches(5.1), COLOR_WHITE, COLOR_BORDER_GRAY)
    
    # 左側表格數據
    headers_inc_ind = ["獲利指標", "113 年度", "112 年度", "增減百分比"]
    rows_inc_ind = [
        ["營業利益 Operating Profit", "10,034,669", "6,231,376", "+61.03%"],
        ["稅前淨利 Profit Before Tax", "15,738,491", "10,483,909", "+50.12%"],
        ["本期淨利 Net Profit", "13,554,170", "9,100,923", "+48.93%"],
        ["每股盈餘 (元) EPS (NT$)", "10.17 元", "6.83 元", "+48.90%"],
        ["其他綜合損益淨額 OCI, net", "(199,284)", "308,934", "-164.51%"],
        ["本期綜合利益總額", "13,354,886", "9,409,857", "+41.92%"]
    ]
    add_styled_table(slide6, left_x, Inches(1.7), left_w, Inches(3.6), headers_inc_ind, rows_inc_ind, widths_inc)
    
    # 右側：原生數據柱狀圖
    add_card(slide6, right_x, Inches(1.7), right_w, Inches(5.1), COLOR_WHITE, COLOR_BORDER_GRAY)
    
    # 建立個體圖表數據 (億元為單位)
    chart_data_ind = CategoryChartData()
    chart_data_ind.categories = ["營業利益", "稅前淨利", "本期淨利"]
    chart_data_ind.add_series("112 年度", (62.3, 104.8, 91.0))
    chart_data_ind.add_series("113 年度", (100.3, 157.4, 135.5))
    
    # 繪製圖表
    chart_shape_ind = slide6.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        right_x + Inches(0.4), Inches(2.2), right_w - Inches(0.8), Inches(4.2), chart_data_ind
    )
    chart_ind = chart_shape_ind.chart
    chart_ind.has_legend = True
    chart_ind.legend.font.name = "微軟正黑體"
    chart_ind.legend.font.size = Pt(10)
    chart_ind.has_title = True
    chart_ind.chart_title.text_frame.text = "個體關鍵獲利指標對比 (新台幣 億元)"
    chart_ind.chart_title.text_frame.paragraphs[0].font.name = "微軟正黑體"
    chart_ind.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart_ind.chart_title.text_frame.paragraphs[0].font.bold = True
    chart_ind.chart_title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY_DARK

    # --------------------------------------------------------------------------
    # 📌 Slide 7: 獲利與每股盈餘 (EPS) 成長亮點專頁 (圖文並茂，視覺極致震撼)
    # --------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLOR_LIGHT_BG)
    add_header(slide7, "每股盈餘 (EPS) 與股本調整亮點分析", "Earnings Per Share & Outstanding Profitability Analysis")
    
    # 左側：大字視覺卡片 (EPS 翻倍)
    left_x = Inches(0.6)
    left_w = Inches(5.8)
    add_card(slide7, left_x, Inches(1.8), left_w, Inches(4.8), COLOR_PRIMARY_DARK)
    
    # 左側文字方塊
    tb_eps = slide7.shapes.add_textbox(left_x + Inches(0.3), Inches(2.0), left_w - Inches(0.6), Inches(4.4))
    tf_eps = tb_eps.text_frame
    tf_eps.word_wrap = True
    
    p_eps_lbl = tf_eps.paragraphs[0]
    p_eps_lbl.text = "113 年度每股盈餘 (EPS)"
    p_eps_lbl.font.name = "微軟正黑體"
    p_eps_lbl.font.size = Pt(16)
    p_eps_lbl.font.color.rgb = COLOR_GOLD
    p_eps_lbl.space_after = Pt(10)
    
    p_eps_num = tf_eps.add_paragraph()
    p_eps_num.text = "10.17 元"
    p_eps_num.font.name = "Arial"
    p_eps_num.font.size = Pt(64)
    p_eps_num.font.bold = True
    p_eps_num.font.color.rgb = COLOR_WHITE
    p_eps_num.space_after = Pt(10)
    
    p_eps_compare = tf_eps.add_paragraph()
    p_eps_compare.text = "相較 112 年度追溯調整後之 6.83 元"
    p_eps_compare.font.name = "微軟正黑體"
    p_eps_compare.font.size = Pt(13)
    p_eps_compare.font.color.rgb = COLOR_TEXT_MUTED
    
    p_eps_growth = tf_eps.add_paragraph()
    p_eps_growth.text = "獲利強勁成長 +48.90%"
    p_eps_growth.font.name = "微軟正黑體"
    p_eps_growth.font.size = Pt(20)
    p_eps_growth.font.bold = True
    p_eps_growth.font.color.rgb = COLOR_GREEN
    p_eps_growth.space_after = Pt(20)
    
    p_eps_foot = tf_eps.add_paragraph()
    p_eps_foot.text = "ℹ️ 註：112 年度之每股盈餘 6.83 元已依 113 年度盈餘轉增資後之發行股數進行追溯調整，以確保比較基礎之一致性。"
    p_eps_foot.font.name = "微軟正黑體"
    p_eps_foot.font.size = Pt(10)
    p_eps_foot.font.color.rgb = COLOR_TEXT_MUTED
    
    # 右側：財務核心成長特點描述
    right_x = Inches(6.8)
    right_w = Inches(5.9)
    add_card(slide7, right_x, Inches(1.8), right_w, Inches(4.8), COLOR_WHITE, COLOR_BORDER_GRAY)
    
    tb_points = slide7.shapes.add_textbox(right_x + Inches(0.4), Inches(2.1), right_w - Inches(0.8), Inches(4.2))
    tf_points = tb_points.text_frame
    tf_points.word_wrap = True
    
    points = [
        ("📈 核心營運利益大增 55.8%+", "合併營業利益由 124 億元躍升至 193.2 億元，反應證券市場交投熱絡與本體核心業務收入的大幅成長。"),
        ("💰 股東權益穩健成長 12.8%", "合併權益總額(分配前)達 1,124.9 億元，保留盈餘更達到 708.6 億元，提供了極佳的營運避險能力與永續經營資本。"),
        ("🔄 盈餘轉增資與股本調整", "113 年順利完成盈餘轉增資，股本擴張的同時 EPS 依然強勁突破兩位數大關，獲利能力卓越。"),
        ("⚡ 個體獲利動能更勝一籌", "證交所本體個體營業利益成長高達 61.03% (由 62.3 億元成長至 100.3 億元)，凸顯交易所本體的高效率獲利體質。")
    ]
    
    for idx, (title, desc) in enumerate(points):
        p_t = tf_points.paragraphs[0] if idx == 0 else tf_points.add_paragraph()
        p_t.text = title
        p_t.font.name = "微軟正黑體"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY_DARK
        
        p_d = tf_points.add_paragraph()
        p_d.text = desc
        p_d.font.name = "微軟正黑體"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_after = Pt(12)
        p_d.line_spacing = 1.25

    # --------------------------------------------------------------------------
    # 📌 Slide 8: 封底 (與封面遙相呼應的暗黑高質感)
    # --------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLOR_DARK_BG)
    
    # 裝飾黃金框
    decor8 = slide8.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.73), Inches(5.9))
    decor8.fill.background()
    decor8.line.color.rgb = COLOR_GOLD
    decor8.line.width = Pt(1.5)
    
    # 結尾文字方塊
    end_box = slide8.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.93), Inches(3.0))
    tf8 = end_box.text_frame
    tf8.word_wrap = True
    
    p8_t = tf8.paragraphs[0]
    p8_t.text = "簡報結束，謝謝聆聽"
    p8_t.alignment = PP_ALIGN.CENTER
    p8_t.font.name = "微軟正黑體"
    p8_t.font.size = Pt(40)
    p8_t.font.bold = True
    p8_t.font.color.rgb = COLOR_WHITE
    p8_t.space_after = Pt(12)
    
    p8_s = tf8.add_paragraph()
    p8_s.text = "TAIWAN STOCK EXCHANGE CORPORATION  |  Concise Financial Report"
    p8_s.alignment = PP_ALIGN.CENTER
    p8_s.font.name = "微軟正黑體"
    p8_s.font.size = Pt(13)
    p8_s.font.bold = True
    p8_s.font.color.rgb = COLOR_GOLD
    
    # ==============================================================================
    # 💾 儲存 PPTX 檔案
    # ==============================================================================
    output_filename = "TWSE_113_Concise_Financial_Report.pptx"
    prs.save(output_filename)
    print(f"🎉 成功生成極致美觀的簡報檔：{output_filename}")

if __name__ == '__main__':
    generate_presentation()
