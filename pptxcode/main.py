import os
import re
import json
import requests
import random
import sys

# 確保在 Windows 環境下輸出 UTF-8 編碼以防止 Emoji 導致 cp950 編碼出錯
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8')
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# =====================================
# ⭐ 新增：載入我們封裝好的雙引擎
# =====================================
from chart_engine import ChartEngine
try:
    from image_engine import ImageEngine
except (ImportError, ModuleNotFoundError):
    ImageEngine = None

# =====================================
# 模型設定
# =====================================
OLLAMA_API = "http://163.13.143.67:1234/api/generate" # 替換為實際連通的遠端伺服器 URL
MODEL_NAME = "gpt-oss:120b"
SAFE_MARGIN = Inches(0.6)

# =====================================
# 🎨 主題庫
# =====================================
THEMES = [
    {
        "name": "BLUE",
        "PRIMARY": RGBColor(37, 99, 235),
        "DARK": RGBColor(30, 41, 59),
        "LIGHT": RGBColor(219, 234, 254),
        "BG": RGBColor(248, 250, 252)
    },
    {
        "name": "BLACK_YELLOW",
        "PRIMARY": RGBColor(234, 179, 8),
        "DARK": RGBColor(0, 0, 0),
        "LIGHT": RGBColor(255, 243, 179),
        "BG": RGBColor(255, 255, 255)
    },
    {
        "name": "RED_PINK",
        "PRIMARY": RGBColor(220, 38, 38),
        "DARK": RGBColor(127, 29, 29),
        "LIGHT": RGBColor(252, 165, 165),
        "BG": RGBColor(255, 240, 240)
    }
]

CURRENT_THEME = random.choice(THEMES)
PRIMARY = CURRENT_THEME["PRIMARY"]
DARK = CURRENT_THEME["DARK"]
LIGHT_BLUE = CURRENT_THEME["LIGHT"]
LIGHT_GRAY = CURRENT_THEME["BG"]
print(f"🎨 使用主題：{CURRENT_THEME['name']}")

# =====================================
# 文字安全與動態字級 (維持原樣)
# =====================================
def enable_textbox_safety(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

def dynamic_font_size(text, base):
    l = len(str(text))
    if l < 20: scale = 1
    elif l < 40: scale = 0.9
    elif l < 80: scale = 0.8
    else: scale = 0.7
    return Pt(base * scale)

def font_color_for_bg(bg_rgb):
    brightness = (0.299*bg_rgb[0] + 0.587*bg_rgb[1] + 0.114*bg_rgb[2])
    return RGBColor(255,255,255) if brightness < 128 else DARK

def generate_short_title(title):
    if len(title) <= 20: return title
    try:
        prompt = f"將以下標題縮短為不超過20字，保持原意完整，不要使用省略號：{title}"
        response = requests.post(OLLAMA_API, json={
            "model": MODEL_NAME,
            "prompt": prompt,
            "stream": False,
            "options": {"temperature": 0.5, "num_predict": 40}
        })
        short_title = response.json().get("response", "").strip()
        if 0 < len(short_title) <= 20: return short_title
    except:
        pass
    return title[:20]

# (為了版面簡潔，我省略了 build_cover_slide 系列、build_toc 等維持原樣的函數，請保留你原本的程式碼)
# ... [保留你原本的 build_cover_slide, build_toc 等函數] ...

def add_title(slide, title_text, sw):
    title_box = slide.shapes.add_textbox(SAFE_MARGIN, Inches(0.6), sw - SAFE_MARGIN*2, Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    enable_textbox_safety(tf)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = dynamic_font_size(title_text, 34)
    p.font.bold = True
    p.font.name = "微軟正黑體"
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.LEFT

def add_card(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    text_color = font_color_for_bg((color[0], color[1], color[2]))
    return shape, text_color

# =====================================
# ⭐ 改造版型：LEFT_TEXT 與 RIGHT_TEXT
# 支援接收雙引擎與視覺化資料
# =====================================
def render_left(slide, data, sw, sh, chart_engine=None, image_engine=None, visual_hint="", chart_infos=None):
    # 1. 左側建立文字卡片
    card_w = sw/2 - SAFE_MARGIN
    card, text_color = add_card(slide, SAFE_MARGIN, Inches(1.6), card_w, sh-Inches(2.3), LIGHT_GRAY)
    tf = card.text_frame
    enable_textbox_safety(tf)
    for s in data.get("sections", []):
        h = tf.add_paragraph()
        h.text = s.get("heading","")
        h.font.bold = True
        h.font.size = dynamic_font_size(h.text,20)
        h.font.color.rgb = text_color
        h.font.name = "微軟正黑體"
        p = tf.add_paragraph()
        p.text = s.get("content","")
        p.font.size = dynamic_font_size(p.text,16)
        p.font.color.rgb = text_color
        p.line_spacing = 1.25
        p.space_after = Pt(12)
        p.font.name = "微軟正黑體"

    # 2. 右側留白區域作為「視覺區」
    vis_x = sw/2 + SAFE_MARGIN/2
    vis_y = Inches(1.6)
    vis_w = sw/2 - SAFE_MARGIN*1.5
    vis_h = sh - Inches(2.3)

    if chart_infos and chart_engine:
        print("   >> 執行繪製圖表任務...")
        chart = chart_infos[0] # 簡化：取第一個圖表繪製
        chart_engine.draw_chart_on_slide(slide, chart["type"], chart["data"], vis_x, vis_y, vis_w, vis_h)
    
    elif visual_hint and image_engine:
        print("   >> 執行生成圖片任務...")
        img_path = f"./temp_images/slide_img_{random.randint(1000,9999)}.png"
        try:
            img_path = image_engine.generate_image(prompt=visual_hint, output_path=img_path)
            slide.shapes.add_picture(img_path, vis_x, vis_y, vis_w, vis_h)
        except Exception as e:
            print(f"❌ 插入圖片失敗: {e}")

def render_right(slide, data, sw, sh, chart_engine=None, image_engine=None, visual_hint="", chart_infos=None):
    # 1. 右側建立文字卡片
    card_x = sw/2 + SAFE_MARGIN/2
    card_w = sw/2 - SAFE_MARGIN*1.5
    card, text_color = add_card(slide, card_x, Inches(1.6), card_w, sh-Inches(2.3), LIGHT_GRAY)
    tf = card.text_frame
    enable_textbox_safety(tf)
    for s in data.get("sections", []):
        h = tf.add_paragraph()
        h.text = s.get("heading","")
        h.font.bold = True
        h.font.size = dynamic_font_size(h.text,20)
        h.font.color.rgb = text_color
        h.font.name = "標楷體"
        p = tf.add_paragraph()
        p.text = s.get("content","")
        p.font.size = dynamic_font_size(p.text,16)
        p.font.color.rgb = text_color
        p.line_spacing = 1.25
        p.space_after = Pt(12)
        p.font.name = "標楷體"

    # 2. 左側留白區域作為「視覺區」
    vis_x = SAFE_MARGIN
    vis_y = Inches(1.6)
    vis_w = sw/2 - SAFE_MARGIN
    vis_h = sh - Inches(2.3)

    if chart_infos and chart_engine:
        print("   >> 執行繪製圖表任務...")
        chart = chart_infos[0]
        chart_engine.draw_chart_on_slide(slide, chart["type"], chart["data"], vis_x, vis_y, vis_w, vis_h)
    
    elif visual_hint and image_engine:
        print("   >> 執行生成圖片任務...")
        img_path = f"./temp_images/slide_img_{random.randint(1000,9999)}.png"
        try:
            img_path = image_engine.generate_image(prompt=visual_hint, output_path=img_path)
            slide.shapes.add_picture(img_path, vis_x, vis_y, vis_w, vis_h)
        except Exception as e:
            print(f"❌ 插入圖片失敗: {e}")

# ... [保留你原本的其他 render 函數: render_three_column, render_comparison 等] ...

ALL_LAYOUTS = [
    "LEFT_TEXT", "RIGHT_TEXT", "THREE_COLUMN", "COMPARISON", "TIMELINE",
    "BIG_NUMBER", "LEFT_NAV_THREE_INSIGHT", "PROCESS_FLOW", "EXECUTIVE_SUMMARY",
    "DIGGING_DEEPER", "DIGGING_DEEPER2"
]

def semantic_layout_selector(data):
    # (維持原樣)
    return "LEFT_TEXT" # 省略內部邏輯，請保留你原來的寫法

LAYOUT_ENGINE = {
    "LEFT_TEXT": render_left,
    "RIGHT_TEXT": render_right,
    # "THREE_COLUMN": render_three_column, 等等... (請保留完整字典)
}

# =====================================
# ⭐ 修改：視覺化建議取得 (引導 LLM 產出好處理的格式)
# =====================================
def get_visual_suggestions(layout, sections):
    prompt = (
        f"針對以下 {layout} 版型的內容，請判斷最適合搭配「數據圖表」還是「情境圖片」。\n"
        f"1. 如果適合數據圖表，請明確回傳包含 PIE, RADAR, LINE 或 BAR 的關鍵字與欄位建議。\n"
        f"2. 如果不適合圖表（純概念或情境），請直接提供一段「全英文」的情境畫面描述（供AI生圖使用），字數在 15 單字以內，不要有其他廢話。\n"
        f"內容: {json.dumps(sections, ensure_ascii=False)}"
    )
    try:
        response = requests.post(
            OLLAMA_API,
            json={
                "model": MODEL_NAME,
                "prompt": prompt,
                "stream": False,
                "options": {"temperature": 0.3, "num_predict": 100}
            }
        )
        return response.json().get("response", "").strip()
    except Exception as e:
        return ""

def smart_layout_decision(data):
    layout = data.get("layout")
    if layout not in ALL_LAYOUTS:
        layout = semantic_layout_selector(data)
    if layout == "DIGGING":
        sec_len = len(data.get("sections", []))
        layout = "DIGGING_DEEPER2" if sec_len >= 3 else "DIGGING_DEEPER"
    return layout

# =====================================
# ⭐ 核心大腦：投影片建置與決策引擎
# =====================================
def build_slide(prs, data, chart_engine, image_engine):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    sh = prs.slide_height

    layout = smart_layout_decision(data)
    
    visual_hint = ""
    chart_infos = []

    # 只有這兩種版型有預留一半的空間
    if layout in ("LEFT_TEXT", "RIGHT_TEXT"):
        visual_hint = get_visual_suggestions(layout, data.get("sections", []))
        if not visual_hint or visual_hint.strip() == "":
            slide_title = data.get("title", "")
            visual_hint = f"A professional high-quality business concept illustration representing {slide_title}, modern corporate design, clean graphics"
        print(f"\n💡 取得視覺建議 ({layout}): {visual_hint}")

        # 1. 先問 ChartEngine 是不是圖表
        chart_types = chart_engine.extract_chart_types(visual_hint)
        
        if chart_types:
            # 是圖表，抽出數據
            slide_text = " ".join([s.get("content", "") for s in data.get("sections", [])])
            chart_infos = chart_engine.analyze_and_get_chart_info(visual_hint, slide_text)
            print(f"📊 判斷為圖表類型: {[c['type'] for c in chart_infos]}")
        else:
            # 沒偵測到圖表關鍵字，當作生圖提示詞
            print(f"🖼️ 判斷為情境圖片需求")

    render_func = LAYOUT_ENGINE.get(layout, render_left)

    # 動態傳遞引擎參數給支援的 Layout
    if layout in ("LEFT_TEXT", "RIGHT_TEXT"):
        render_func(slide, data, sw, sh, chart_engine, image_engine, visual_hint, chart_infos)
    else:
        render_func(slide, data, sw, sh)

    if layout not in ["LEFT_NAV_THREE_INSIGHT", "PROCESS_FLOW", "EXECUTIVE_SUMMARY", "DIGGING_DEEPER", "DIGGING_DEEPER2"]:
        add_title(slide, data.get("title", ""), sw)

# =====================================
# 建立 PPT (加入雙引擎參數)
# =====================================
def create_ppt(slides, chart_engine, image_engine):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    if slides:
        cover_slide = slides[0]
        content_slides = [s for s in slides[1:] if s.get("title", "").strip() not in ["目錄", "contents", "table of contents"]]
        slides_clean = [cover_slide] + content_slides

        build_cover(prs, slides_clean) 
        build_toc_slide(prs, slides_clean)

        for s in slides_clean[1:]:
            build_slide(prs, s, chart_engine, image_engine)

    import shutil
    try:
        prs.save("Smart_Presentation_With_Visuals.pptx")
        shutil.copy("Smart_Presentation_With_Visuals.pptx", "C:/Users/user/Desktop/Smart_Presentation_With_Visuals.pptx")
        print("\n🎉 完成並已複製至桌面: C:/Users/user/Desktop/Smart_Presentation_With_Visuals.pptx")
    except Exception as e:
        print(f"⚠️ 儲存或複製簡報失敗: {e}")

def set_slide_background_main(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def build_cover(prs, slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background_main(slide, DARK)
    
    # 標題
    title_box = slide.shapes.add_textbox(SAFE_MARGIN, Inches(2.2), prs.slide_width - SAFE_MARGIN*2, Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    enable_textbox_safety(tf)
    
    p = tf.paragraphs[0]
    p.text = slides[0].get("title", "113年度簡明財務報告")
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = "微軟正黑體"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    p2 = tf.add_paragraph()
    p2.text = "Taiwan Stock Exchange Concise Financial Report"
    p2.font.size = Pt(18)
    p2.font.name = "微軟正黑體"
    p2.font.color.rgb = PRIMARY
    p2.alignment = PP_ALIGN.CENTER

def build_toc_slide(prs, slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background_main(slide, LIGHT_GRAY)
    add_title(slide, "目錄", prs.slide_width)
    
    # 目錄列表
    box = slide.shapes.add_textbox(SAFE_MARGIN, Inches(1.8), prs.slide_width - SAFE_MARGIN*2, Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    enable_textbox_safety(tf)
    
    for idx, s in enumerate(slides[1:]):
        p = tf.add_paragraph()
        p.text = f"{idx+1:02d}. {s.get('title', '')}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = "微軟正黑體"
        p.font.color.rgb = DARK
        p.space_after = Pt(12)

def get_slide_structure(text):
    prompt = f"""
    請閱讀以下公司的財務報告內容，並將其分析、摘要，轉化為一份適合用於高階商業簡報的結構。
    總共需要規劃 5-7 張投影片。請直接回傳 JSON 格式陣列，不要有任何 Markdown 的 ```json 語法標記或其它多餘文字。
    
    每個投影片物件必須包含以下欄位：
    - "title": 該頁投影片標題 (繁體中文)
    - "layout": 必須是 "LEFT_TEXT" 或 "RIGHT_TEXT" 之一
    - "sections": 陣列，每個元素為一個區段字典，含 "heading" (小標題) 與 "content" (詳細描述，約 50 字內，不要太長以防排版溢出)
    
    財務報告內容：
    {text}
    """
    try:
        response = requests.post(OLLAMA_API, json={
            "model": MODEL_NAME,
            "prompt": prompt,
            "stream": False,
            "options": {"temperature": 0.3}
        }, timeout=120)
        return response.json().get("response", "").strip()
    except Exception as e:
        print(f"❌ 呼叫 Ollama 失敗: {e}")
        return ""

def parse_slides(raw_text):
    if not raw_text:
        return []
    try:
        # 清理並尋找 JSON 區塊 [ ... ]
        match = re.search(r'\[\s*\{.*\}\s*\]', raw_text, re.DOTALL)
        if match:
            return json.loads(match.group(0))
        # 嘗試直接解析
        return json.loads(raw_text.strip())
    except Exception as e:
        print(f"❌ JSON 解析投影片結構失敗: {e}")
        # 提供一套預設的投影片結構以防萬一
        return [
            {
                "title": "臺灣證券交易所 113年簡明財務報告",
                "layout": "LEFT_TEXT",
                "sections": [{"heading": "財務報告概述", "content": "本簡報呈現臺灣證券交易所 113 年度財務成果與合併報表績效分析。"}]
            }
        ]

# (避免程式報錯，若你原本有這函數請保留)
def ensure_all_layouts(slides):
    return slides

# =====================================
# 主程式
# =====================================
if __name__ == "__main__":
    WORD_PATH = r"C:/Users/user/Desktop/txttoimg/txttoimg/0_H_26624 (1).docx"

    # 確保暫存圖片的資料夾存在
    os.makedirs("./temp_images", exist_ok=True)
    
    print("啟動雙視覺引擎...")
    # 啟動時只載入一次模型，避免生成每一頁都要等
    my_chart_engine = ChartEngine(ollama_url=OLLAMA_API, ollama_model=MODEL_NAME)
    
    my_image_engine = None
    try:
        # 💡 將原本的本地路徑，改回官方的 Repository 名稱
        my_image_engine = ImageEngine(model_id_or_path="KBlueLeaf/kohaku-v2.1")
    except Exception as e:
        print(f"⚠️ 警告：無法初始化 StreamDiffusion 圖片生成引擎 ({e})")
        print("💡 系統將以「文字 + 原生財務圖表」之優雅商業版型繼續生成簡報。")

    if os.path.exists(WORD_PATH):
        doc = Document(WORD_PATH)
        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        print("🚀 AI 分析中並自動配置最適版型...")
        
        raw = get_slide_structure(text) # AI 解析器
        slides = parse_slides(raw)
        slides = ensure_all_layouts(slides)
        
        if slides:
            # 執行建立並傳入雙引擎
            create_ppt(slides, my_chart_engine, my_image_engine)
        else:
            print("❌ AI 未成功產生 JSON")
    else:
        print(f"❌ 找不到檔案: {WORD_PATH}")
